# -*- coding: utf-8 -*-
from pathlib import Path
import sys

sys.path.insert(0, str(Path(".codex_deps").resolve()))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "基于多模态学习的药物分子性质预测研究_答辩PPT.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT = "Microsoft YaHei"
FONT_EN = "Aptos"

INK = "17212B"
GREEN = "174C43"
TEAL = "0E9F9A"
MINT = "DDF7F2"
AMBER = "F2B84B"
PAPER = "F7F3EA"
LIGHT = "F9FBF8"
GRAY = "5E6876"
MID = "D7E6DE"
WHITE = "FFFFFF"
RED = "B33A3A"


def rgb(hex_color):
    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def inch(value):
    return Inches(value)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align="left",
             valign="top", font=FONT, margin=0.05, line_spacing=1.05):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(margin)
    tf.margin_bottom = inch(margin)
    tf.word_wrap = True
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)

    for idx, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(align, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_section_title(slide, title, subtitle=None, section=None, dark=False):
    color = WHITE if dark else INK
    if section:
        add_pill(slide, section, 0.62, 0.42, 1.35, 0.34, TEAL if dark else GREEN, WHITE, 10)
    add_text(slide, title, 0.62, 0.78, 8.6, 0.55, size=28, color=color, bold=True, margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.34, 8.8, 0.38, size=11.5, color=MID if dark else GRAY, margin=0)


def add_footer(slide, page, dark=False):
    color = "B9C7C0" if dark else "7B858E"
    add_text(slide, "基于多模态学习的药物分子性质预测研究", 0.62, 7.08, 5.6, 0.22,
             size=8.5, color=color, margin=0)
    add_text(slide, f"{page:02d}", 12.25, 7.03, 0.55, 0.28, size=9, color=color,
             align="right", margin=0)


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, radius=False):
    shp = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_pill(slide, text, x, y, w, h, fill, color=WHITE, size=11, bold=True):
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, None)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = inch(0.06)
    tf.margin_right = inch(0.06)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shp


def add_card(slide, title, body, x, y, w, h, accent=TEAL, fill=WHITE, number=None):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, "DDE6E0")
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent, None)
    if number:
        add_pill(slide, number, x + 0.24, y + 0.24, 0.42, 0.36, accent, WHITE, 10)
        tx = x + 0.78
        tw = w - 1.0
    else:
        tx = x + 0.28
        tw = w - 0.45
    add_text(slide, title, tx, y + 0.20, tw, 0.33, size=15.5, color=INK, bold=True, margin=0)
    add_text(slide, body, x + 0.30, y + 0.68, w - 0.55, h - 0.82, size=11.3, color=GRAY,
             margin=0, line_spacing=1.08)


def add_metric(slide, value, label, x, y, w, h, color=TEAL):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, WHITE, "DCE5E1")
    add_text(slide, value, x + 0.08, y + 0.12, w - 0.16, 0.42, size=23, color=color,
             bold=True, align="center", margin=0)
    add_text(slide, label, x + 0.12, y + 0.65, w - 0.24, h - 0.74, size=9.5, color=GRAY,
             align="center", margin=0)


def add_image_contain(slide, path, x, y, w, h, caption=None, border=True, pad=0.09, bg=WHITE):
    path = Path(path)
    if not path.is_absolute():
        path = ROOT / path
    if bg:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, bg, "DDE6E0" if border else None)
    img = Image.open(path)
    iw, ih = img.size
    cap_h = 0.28 if caption else 0
    max_w = w - 2 * pad
    max_h = h - 2 * pad - cap_h
    ratio = min(max_w / iw, max_h / ih)
    draw_w = iw * ratio
    draw_h = ih * ratio
    px = x + (w - draw_w) / 2
    py = y + pad + (max_h - draw_h) / 2
    slide.shapes.add_picture(str(path), inch(px), inch(py), width=inch(draw_w), height=inch(draw_h))
    if caption:
        add_text(slide, caption, x + 0.16, y + h - 0.30, w - 0.32, 0.20, size=8.5,
                 color=GRAY, align="center", margin=0)


def add_image_bleed(slide, path, x, y, w, h):
    path = Path(path)
    if not path.is_absolute():
        path = ROOT / path
    img = Image.open(path)
    iw, ih = img.size
    ratio = max(w / iw, h / ih)
    draw_w = iw * ratio
    draw_h = ih * ratio
    px = x + (w - draw_w) / 2
    py = y + (h - draw_h) / 2
    slide.shapes.add_picture(str(path), inch(px), inch(py), width=inch(draw_w), height=inch(draw_h))


def add_flow(slide, steps, x, y, w, h, colors):
    gap = 0.22
    item_w = (w - gap * (len(steps) - 1)) / len(steps)
    for i, (title, desc) in enumerate(steps):
        sx = x + i * (item_w + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, sx, y, item_w, h, colors[i % len(colors)], None)
        add_text(slide, f"{i + 1}", sx + 0.12, y + 0.12, 0.35, 0.3, size=13, color=WHITE,
                 bold=True, margin=0)
        add_text(slide, title, sx + 0.18, y + 0.55, item_w - 0.36, 0.33, size=13.3,
                 color=WHITE, bold=True, align="center", margin=0)
        add_text(slide, desc, sx + 0.20, y + 1.02, item_w - 0.40, h - 1.15, size=9.2,
                 color="EEF8F4", align="center", margin=0, line_spacing=1.08)
        if i < len(steps) - 1:
            add_text(slide, "→", sx + item_w + 0.02, y + h / 2 - 0.18, 0.18, 0.3,
                     size=16, color=GRAY, bold=True, align="center", margin=0)


def add_timeline_item(slide, num, title, body, x, y, w):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, 0.44, 0.44, TEAL if num % 2 else AMBER, None)
    add_text(slide, str(num), x, y + 0.075, 0.44, 0.18, size=10.5, color=WHITE,
             bold=True, align="center", margin=0)
    add_text(slide, title, x + 0.68, y - 0.02, w - 0.68, 0.26, size=14.5, color=INK,
             bold=True, margin=0)
    add_text(slide, body, x + 0.68, y + 0.30, w - 0.68, 0.38, size=10.5, color=GRAY,
             margin=0)


def make_deck():
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_background(slide, INK)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, INK, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.22, 7.5, TEAL, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.22, 0, 0.08, 7.5, AMBER, None)
    add_pill(slide, "专业硕士学位论文答辩", 0.82, 0.74, 1.9, 0.38, GREEN, WHITE, 10.5)
    add_text(slide, "基于多模态学习的\n药物分子性质预测研究", 0.82, 1.45, 6.0, 1.55,
             size=31, color=WHITE, bold=True, margin=0, line_spacing=0.95)
    add_text(slide, "A Multimodal Learning Approach to Drug Molecule Property Prediction",
             0.84, 3.24, 6.4, 0.30, size=11.5, color=MID, font=FONT_EN, margin=0)
    add_text(slide, "研究生：范帅尧    专业：计算机技术\n校内导师：牛云云    产业导师：张永虹\n中国地质大学（北京） · 2026 年 04 月",
             0.84, 5.62, 5.7, 0.86, size=12, color="DDE7E2", margin=0, line_spacing=1.1)
    add_image_contain(slide, "image/第二章相关理论基础/图2-1 阿司匹林的多模态化学表征.png",
                      7.35, 0.95, 5.15, 5.45, caption="分子的图结构、SMILES、指纹与文本语义可共同构成多模态表征", pad=0.18)
    add_footer(slide, 1, dark=True)

    # 2. Outline
    slide = prs.slides.add_slide(blank)
    set_background(slide, LIGHT)
    add_section_title(slide, "汇报提纲", "围绕“问题—方法—实验—结论”的答辩主线展开", "Outline")
    items = [
        ("研究背景与问题", "药物早筛对分子性质预测的需求，以及现有方法的局限"),
        ("相关理论基础", "分子多模态表征、图神经网络、Transformer、LoRA 与思维链"),
        ("方法一：CoT-CMP", "思维链增强文本推理 + 通信式图消息传递 + 图文交叉注意力"),
        ("方法二：KGAMA", "知识重写、四模态表征、图中心对齐和不确定性加权"),
        ("总结与展望", "核心创新、实验验证、局限性与未来工作"),
    ]
    for i, item in enumerate(items, 1):
        add_timeline_item(slide, i, item[0], item[1], 1.0, 1.63 + (i - 1) * 0.92, 6.2)
    add_image_contain(slide, "image/第四章/KGAMA模型架构.png", 7.35, 1.35, 5.1, 4.95,
                      caption="答辩核心围绕两个多模态分子性质预测框架展开", pad=0.14)
    add_footer(slide, 2)

    # 3. Background
    slide = prs.slides.add_slide(blank)
    set_background(slide, PAPER)
    add_section_title(slide, "研究背景：药物早筛需要更可靠的分子性质预测",
                      "在候选药物数量巨大、实验验证昂贵的场景中，计算预测可以前置风险识别", "Background")
    add_text(slide, "现代创新药研发周期长、成本高，早期筛选阶段需要对 ADMET、毒性、溶解度、亲脂性等性质进行快速评估。",
             0.82, 1.46, 5.3, 0.72, size=15.5, color=INK, bold=True, margin=0)
    add_card(slide, "核心目标", "把复杂分子结构与语义知识映射为稳定表示，并服务于分类与回归两类性质预测任务。", 0.82, 2.45, 3.65, 1.24, TEAL)
    add_card(slide, "关键挑战", "分子既包含拓扑结构，也包含序列、指纹、理化指标和文献描述，单一模态很难完整表达。", 4.75, 2.45, 3.65, 1.24, AMBER)
    steps = [
        ("候选库", "百万级化合物"),
        ("性质预测", "ADMET / 毒性"),
        ("风险筛除", "降低后期失败率"),
        ("实验验证", "聚焦高潜力分子"),
    ]
    add_flow(slide, steps, 0.82, 4.42, 7.42, 1.42, [GREEN, TEAL, AMBER, INK])
    add_image_contain(slide, "image/第二章相关理论基础/图2-2 消息传递机制示意图.png",
                      8.72, 1.45, 3.72, 4.85, caption="图神经网络为分子结构建模提供基础，但仍存在局限", pad=0.16)
    add_footer(slide, 3)

    # 4. Gaps and idea
    slide = prs.slides.add_slide(blank)
    set_background(slide, LIGHT)
    add_section_title(slide, "当前研究不足与本文整体思路", "本文从“可解释推理”和“稳健对齐”两个角度展开", "Problem")
    add_card(slide, "拓扑感知受限", "传统 GNN 依赖局部节点聚合，化学键动态交互不足，深层传播容易带来信息冗余与回流。", 0.82, 1.52, 3.65, 1.58, RED, number="01")
    add_card(slide, "LLM 推理难解释", "大语言模型直接预测复杂 ADMET 性质时，缺少中间逻辑约束，容易产生幻觉且不利于审计。", 4.82, 1.52, 3.65, 1.58, AMBER, number="02")
    add_card(slide, "文本噪声与对齐困难", "公共医学文本存在冗余和质量差异，跨模态对比学习容易受到语义冲突与负迁移影响。", 8.82, 1.52, 3.65, 1.58, TEAL, number="03")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.05, 4.10, 5.15, 1.28, INK, None)
    add_text(slide, "CoT-CMP", 1.38, 4.31, 1.35, 0.32, size=17, color=WHITE, bold=True, margin=0)
    add_text(slide, "思维链增强文本端可解释性，通信式消息传递提升图结构表征能力。", 1.38, 4.76, 4.38, 0.38, size=11.5, color="E6EFEA", margin=0)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.10, 4.10, 5.15, 1.28, GREEN, None)
    add_text(slide, "KGAMA", 7.43, 4.31, 1.35, 0.32, size=17, color=WHITE, bold=True, margin=0)
    add_text(slide, "知识重写提高文本质量，图中心对齐与不确定性加权稳定四模态融合。", 7.43, 4.76, 4.38, 0.38, size=11.5, color="E6EFEA", margin=0)
    add_text(slide, "两条研究线索共同服务于分子性质预测的准确性、可解释性与泛化能力。", 1.05, 5.86, 11.2, 0.38,
             size=15, color=INK, bold=True, align="center", margin=0)
    add_footer(slide, 4)

    # 5. Foundations
    slide = prs.slides.add_slide(blank)
    set_background(slide, PAPER)
    add_section_title(slide, "相关理论基础：多模态分子表征与深度模型",
                      "本文以图、序列、指纹、文本四类信息构建分子表示", "Foundation")
    add_image_contain(slide, "image/第二章相关理论基础/图2-1 阿司匹林的多模态化学表征.png",
                      0.82, 1.34, 5.95, 2.78, caption="同一分子可从多种模态获得互补信息", pad=0.14)
    add_image_contain(slide, "image/第二章相关理论基础/图2-4 Transformer架构.png",
                      7.10, 1.34, 2.68, 2.78, caption="Transformer 建模长程依赖", pad=0.13)
    add_image_contain(slide, "image/第二章相关理论基础/LoRA技术原理示意图.png",
                      10.06, 1.34, 2.38, 2.78, caption="LoRA 轻量化适配 LLM", pad=0.13)
    rows = [
        ("二维图结构", "显式保留原子与化学键连接关系，是结构建模的核心。"),
        ("SMILES 序列", "以线性符号序列承载化学语法，适合 Transformer 编码。"),
        ("分子指纹", "统计特定子结构和官能团，补充全局结构先验。"),
        ("自然语言文本", "承载药理、毒理与领域知识，但需要去噪和约束。"),
    ]
    for idx, (t, b) in enumerate(rows):
        add_card(slide, t, b, 0.82 + (idx % 2) * 6.03, 4.48 + (idx // 2) * 0.95, 5.58, 0.70,
                 [TEAL, GREEN, AMBER, INK][idx])
    add_footer(slide, 5)

    # 6. CoT-CMP architecture
    slide = prs.slides.add_slide(blank)
    set_background(slide, LIGHT)
    add_section_title(slide, "方法一：CoT-CMP 双塔预测框架", "面向 ADMET 任务，融合显式药理推理文本与交互式图结构表示", "Method 1")
    add_image_contain(slide, "image/第三章/图3-1CoTCMP架构图.png", 0.82, 1.30, 11.70, 3.55,
                      caption="CoT-CMP：文本编码、图编码与融合预测三部分组成", pad=0.13)
    add_card(slide, "文本端", "LLaMA-3.1-8B + LoRA + CoT 提示，引导生成“结构识别—理化分析—性质推断”的显式推理文本。", 0.92, 5.20, 3.72, 1.02, TEAL)
    add_card(slide, "图结构端", "CMPNN 通过原子与化学键通信核建模动态交互，缓解传统消息传递的拓扑感知不足。", 4.93, 5.20, 3.72, 1.02, GREEN)
    add_card(slide, "融合预测", "以交叉注意力机制对齐图特征与文本特征，实现微观拓扑与宏观药理语义的自适应融合。", 8.94, 5.20, 3.72, 1.02, AMBER)
    add_footer(slide, 6)

    # 7. CoT-CMP mechanisms
    slide = prs.slides.add_slide(blank)
    set_background(slide, PAPER)
    add_section_title(slide, "CoT-CMP 核心机制", "将大模型的推理能力约束为可用的分子性质预测特征", "Method 1")
    add_image_contain(slide, "table/第三章/提示词模板.png", 0.82, 1.32, 3.72, 3.88,
                      caption="CoT 提示词模板：显式约束推理路径", pad=0.12)
    add_image_contain(slide, "image/第二章相关理论基础/图2-6 输入-输出提示与思维链提示对比示意图.png",
                      4.82, 1.32, 2.62, 3.88, caption="从直接输出转向分步推理", pad=0.12)
    add_image_contain(slide, "image/第二章相关理论基础/图2-2 消息传递机制示意图.png",
                      7.72, 1.32, 4.72, 2.06, caption="消息传递机制", pad=0.12)
    add_card(slide, "为什么需要 CoT", "直接预测容易缺少药理逻辑；分步推理把结构片段、理化性质和目标性质之间的关系显式化。", 7.72, 3.75, 4.72, 0.78, TEAL)
    add_card(slide, "为什么需要 CMPNN", "将边特征纳入通信过程，使化学键不再只是辅助输入，而是共同参与分子结构建模。", 7.72, 4.78, 4.72, 0.78, GREEN)
    add_card(slide, "为什么需要交叉注意力", "简单拼接无法刻画图结构与文本语义的对应关系；交叉注意力用于动态选择互补信息。", 7.72, 5.81, 4.72, 0.78, AMBER)
    add_footer(slide, 7)

    # 8. CoT-CMP experiments
    slide = prs.slides.add_slide(blank)
    set_background(slide, LIGHT)
    add_section_title(slide, "CoT-CMP 实验设计", "在 TDC ADMET 基准组上验证分类与回归预测能力", "Experiment")
    add_image_contain(slide, "table/第三章/ADMET基准数据集.png", 0.82, 1.28, 4.55, 4.95,
                      caption="ADMET 基准数据集覆盖吸收、分布、代谢、排泄与毒性", pad=0.12)
    add_image_contain(slide, "image/第三章/随机划分与骨架划分在PPBR_AZ数据集上分布的对比.png",
                      5.68, 1.28, 6.78, 2.45, caption="随机划分与骨架划分体现不同泛化评估难度", pad=0.12)
    add_card(slide, "实验设置", "LLaMA-3.1-8B-Instruct 作为文本主干；LoRA 参数 r=16、缩放系数 32；图编码器采用 3 层 CMPNN；4bit 量化降低显存占用。", 5.68, 4.05, 3.22, 1.28, TEAL)
    add_card(slide, "评价指标", "分类任务采用 AUROC/AUPRC；回归任务采用 MAE 与 Spearman，更贴近药物筛选中的排序需求。", 9.24, 4.05, 3.22, 1.28, GREEN)
    add_metric(slide, "22", "ADMET 评测数据集", 5.68, 5.64, 1.55, 0.76, TEAL)
    add_metric(slide, "200", "最大训练轮次", 7.48, 5.64, 1.55, 0.76, AMBER)
    add_metric(slide, "15", "早停 patience", 9.28, 5.64, 1.55, 0.76, GREEN)
    add_metric(slide, "4090", "24GB GPU", 11.08, 5.64, 1.38, 0.76, INK)
    add_footer(slide, 8)

    # 9. CoT-CMP results
    slide = prs.slides.add_slide(blank)
    set_background(slide, PAPER)
    add_section_title(slide, "CoT-CMP 对比实验结果", "在 22 个 ADMET 数据集上整体优于经典图模型与序列模型", "Result")
    add_image_contain(slide, "table/第三章/CoT-CMP在ADMET基准组上的对比实验结果.png",
                      0.82, 1.24, 6.18, 5.22, caption="表 3-3：CoT-CMP 在 ADMET 基准组上的对比实验结果", pad=0.10)
    add_metric(slide, "17/22", "取得最优结果", 7.38, 1.42, 1.75, 0.92, TEAL)
    add_metric(slide, "5/22", "取得次优结果", 9.38, 1.42, 1.75, 0.92, AMBER)
    add_metric(slide, "4/4", "毒性任务全面领先", 11.38, 1.42, 1.35, 0.92, RED)
    add_card(slide, "毒性预测表现突出", "hERG、Ames、DILI 等毒性相关任务均达到最高预测精度，支持候选药物早期风险评估。", 7.38, 2.78, 5.35, 1.02, RED)
    add_card(slide, "图文融合带来鲁棒表征", "与 AttentiveFP、NeuralFP 等图神经网络相比，CoT 增强文本与交叉注意力进一步聚合长程依赖和局部官能团特征。", 7.38, 4.02, 5.35, 1.20, TEAL)
    add_card(slide, "回归任务同样有效", "在 PPBR、LD50、Caco2、AqSol 等连续值预测中，模型能降低误差并保持排序能力。", 7.38, 5.47, 5.35, 0.98, GREEN)
    add_footer(slide, 9)

    # 10. CoT-CMP ablation and explanation
    slide = prs.slides.add_slide(blank)
    set_background(slide, LIGHT)
    add_section_title(slide, "CoT-CMP 消融与可解释性分析", "模块贡献与可视化结果共同验证模型设计有效性", "Result")
    add_image_contain(slide, "table/第三章/CoT-CMP在ADMET基准组上部分分类任务消融实验结果.png",
                      0.82, 1.25, 5.42, 1.88, caption="分类任务消融：完整模型整体最优", pad=0.12)
    add_image_contain(slide, "image/第三章/氯磺丙脲在DILI任务中的双模态可视化分析结果.png",
                      0.82, 3.48, 5.42, 2.78, caption="DILI 案例：结构热点与毒理关键词相互印证", pad=0.12)
    add_image_contain(slide, "image/第三章/CMPNN模型层数对分类任务性能的影响.png",
                      6.72, 1.25, 5.72, 2.58, caption="参数敏感性：3 层 CMPNN 在全局感受野与过拟合之间取得平衡", pad=0.12)
    add_card(slide, "消融结论", "移除 LoRA、CoT、LLM 或交叉注意力均会带来不同程度下降，说明文本增强、任务适配与跨模态融合之间存在协同增益。", 6.72, 4.22, 5.72, 0.92, TEAL)
    add_card(slide, "解释结论", "氯磺丙脲案例中，模型关注对氯苯基、磺酰脲骨架以及 hepatotoxicity、DILI 等关键词，与药物性肝损伤判断直觉一致。", 6.72, 5.42, 5.72, 0.92, GREEN)
    add_footer(slide, 10)

    # 11. KGAMA architecture
    slide = prs.slides.add_slide(blank)
    set_background(slide, PAPER)
    add_section_title(slide, "方法二：KGAMA 知识增强与图对齐框架", "面向 MoleculeNet，构建图、序列、指纹、文本四模态表征体系", "Method 2")
    add_image_contain(slide, "image/第四章/KGAMA模型架构.png", 0.82, 1.20, 7.50, 4.95,
                      caption="KGAMA：Graph-centered Alignment + Knowledge-enhanced Multimodal Architecture", pad=0.12)
    add_card(slide, "四模态输入", "分子图、SMILES 序列、分子指纹和自然语言文本共同刻画结构、语法、子结构统计与药理语义。", 8.72, 1.32, 3.72, 1.10, TEAL)
    add_card(slide, "图中心对齐", "以分子图作为语义锚点，避免全排列式对齐带来的冗余冲突，使各辅助模态围绕稳定结构表征收敛。", 8.72, 2.72, 3.72, 1.16, GREEN)
    add_card(slide, "自适应加权", "基于同方差不确定性调节对比学习损失权重，降低噪声模态对总梯度的干扰。", 8.72, 4.20, 3.72, 1.10, AMBER)
    add_card(slide, "知识增强", "通过 RDKit 理化指标与 LLM 重写提升文本信息密度，缓解公共医学文本冗余与噪声。", 8.72, 5.60, 3.72, 0.88, RED)
    add_footer(slide, 11)

    # 12. KGAMA knowledge rewrite
    slide = prs.slides.add_slide(blank)
    set_background(slide, LIGHT)
    add_section_title(slide, "KGAMA：知识重写与模态构建", "把低质量原始文本转化为可被模型利用的高密度药理语义", "Method 2")
    add_flow(slide, [
        ("PubChem", "原始分子描述"),
        ("RDKit", "11 项理化指标"),
        ("LLM 重写", "去噪与知识注入"),
        ("SciBERT", "语义编码"),
    ], 0.82, 1.42, 6.18, 1.32, [INK, GREEN, TEAL, AMBER])
    add_image_contain(slide, "image/第四章/药物分子原始检索文本与重写文本的对比示例.png",
                      7.36, 1.16, 4.95, 5.48, caption="知识重写示例：压缩冗余文本并显式加入理化属性", pad=0.12)
    add_card(slide, "文本端知识注入", "利用 RDKit 补充 LogP、TPSA 等关键理化属性，让文本模态从描述性语料转为化学约束更强的专家语义。", 0.82, 3.22, 2.98, 1.15, TEAL)
    add_card(slide, "指纹端结构补充", "AtomPairs、MACCS、Morgan 三类指纹覆盖不同子结构统计信息，补充图编码可能遗漏的官能团数量特征。", 4.10, 3.22, 2.98, 1.15, GREEN)
    add_image_contain(slide, "table/第四章/三种典型分子指纹特征的比较分析.png",
                      0.82, 4.75, 2.98, 1.18, caption="指纹特征比较", pad=0.10)
    add_image_contain(slide, "table/第四章/分子理化性质指标及其含义.png",
                      4.10, 4.75, 2.98, 1.18, caption="理化性质指标", pad=0.10)
    add_footer(slide, 12)

    # 13. KGAMA comparison
    slide = prs.slides.add_slide(blank)
    set_background(slide, PAPER)
    add_section_title(slide, "KGAMA 对比实验结果", "在 MoleculeNet 九个基准数据集上验证分类与回归能力", "Result")
    add_image_contain(slide, "table/第四章/KGAMA在MoleculeNet分类任务上对比实验结果.png",
                      0.82, 1.22, 6.18, 4.82, caption="表 4-5：KGAMA 在 MoleculeNet 分类任务上的对比实验结果", pad=0.10)
    add_metric(slide, "0.939", "ClinTox AUROC", 7.38, 1.34, 1.75, 0.88, RED)
    add_metric(slide, "0.930", "BBBP AUROC", 9.38, 1.34, 1.75, 0.88, TEAL)
    add_metric(slide, "0.669", "SIDER AUROC", 11.38, 1.34, 1.35, 0.88, GREEN)
    add_card(slide, "分类任务", "KGAMA 在 BACE、BBBP、ClinTox、Tox21、ToxCast、SIDER 六个分类数据集上均取得稳定优势。", 7.38, 2.62, 5.35, 0.96, TEAL)
    add_card(slide, "回归任务", "ESOL、FreeSolv、Lipo 三个连续性质任务中分别达到 0.831、1.512、0.655 的 RMSE 表现。", 7.38, 3.84, 5.35, 0.96, GREEN)
    add_card(slide, "优势来源", "Graphormer 捕获全局拓扑，知识重写提升文本质量，不确定性加权增强复杂任务下的泛化稳定性。", 7.38, 5.06, 5.35, 0.96, AMBER)
    add_footer(slide, 13)

    # 14. KGAMA ablation and attention
    slide = prs.slides.add_slide(blank)
    set_background(slide, LIGHT)
    add_section_title(slide, "KGAMA 消融与可解释性分析", "注意力权重揭示不同任务对不同模态的依赖差异", "Result")
    add_image_contain(slide, "image/第四章/KGAMA在MoleculeNet基准组上分类任务消融实验结果.png",
                      0.82, 1.23, 5.55, 2.38, caption="分类任务消融：模态逐步加入带来性能递增", pad=0.12)
    add_image_contain(slide, "image/第四章/KGAMA在MoleculeNet基准组上回归任务消融实验结果.png",
                      0.82, 3.92, 5.55, 2.22, caption="回归任务消融：完整模型误差最低", pad=0.12)
    add_image_contain(slide, "image/第四章/不同分子模态在各下游任务中的注意力权重分布.png",
                      6.78, 1.23, 5.62, 2.96, caption="注意力权重：不同任务自动选择不同模态", pad=0.12)
    add_card(slide, "模态贡献", "Graph 在 ESOL、FreeSolv 等物理性质任务中权重最高；Fingerprint 在 Lipo 中更关键；Text 在 SIDER、ClinTox、Tox21 等药理毒理任务中显著跃升。", 6.78, 4.54, 5.62, 0.96, TEAL)
    add_card(slide, "融合机制贡献", "完整四模态输入加交叉注意力优于静态拼接，说明模型不是简单叠加信息，而是根据任务动态调配模态权重。", 6.78, 5.72, 5.62, 0.82, GREEN)
    add_footer(slide, 14)

    # 15. Conclusions
    slide = prs.slides.add_slide(blank)
    set_background(slide, PAPER)
    add_section_title(slide, "工作总结与创新点", "两项方法共同提升多模态分子性质预测的准确性、可解释性与稳定性", "Conclusion")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.90, 1.40, 5.55, 4.45, INK, None)
    add_text(slide, "CoT-CMP", 1.25, 1.78, 2.2, 0.38, size=20, color=WHITE, bold=True, margin=0)
    add_text(slide, "面向 ADMET 任务的图文双塔框架", 1.25, 2.22, 4.35, 0.28, size=11.5, color=MID, margin=0)
    add_text(slide, "1. CoT 提示把 LLM 推理过程显式化\n2. LoRA 实现面向分子任务的高效适配\n3. CMPNN 强化原子—化学键通信\n4. 交叉注意力融合文本逻辑与图结构",
             1.25, 2.85, 4.65, 1.35, size=14, color="EFF6F2", margin=0, line_spacing=1.15)
    add_metric(slide, "17/22", "ADMET 数据集最优", 1.25, 4.82, 1.85, 0.78, TEAL)
    add_metric(slide, "4/4", "毒性任务领先", 3.35, 4.82, 1.85, 0.78, RED)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.90, 1.40, 5.55, 4.45, GREEN, None)
    add_text(slide, "KGAMA", 7.25, 1.78, 2.2, 0.38, size=20, color=WHITE, bold=True, margin=0)
    add_text(slide, "面向 MoleculeNet 的知识增强四模态框架", 7.25, 2.22, 4.45, 0.28, size=11.5, color="E6EFEA", margin=0)
    add_text(slide, "1. LLM 知识重写提升文本质量\n2. RDKit 理化指标提供显式约束\n3. 图中心对齐减少多模态冲突\n4. 不确定性加权提升训练稳定性",
             7.25, 2.85, 4.65, 1.35, size=14, color="EFF6F2", margin=0, line_spacing=1.15)
    add_metric(slide, "9", "MoleculeNet 数据集", 7.25, 4.82, 1.85, 0.78, AMBER)
    add_metric(slide, "4", "分子信息模态", 9.35, 4.82, 1.85, 0.78, TEAL)
    add_text(slide, "总体结论：多模态并非简单堆叠，关键在于让每种模态“有知识、有约束、能对齐、可解释”。",
             1.20, 6.42, 11.0, 0.38, size=15.5, color=INK, bold=True, align="center", margin=0)
    add_footer(slide, 15)

    # 16. Outlook
    slide = prs.slides.add_slide(blank)
    set_background(slide, INK)
    add_section_title(slide, "工作展望", "进一步面向真实药物发现流程提升建模能力与验证闭环", "Outlook", dark=True)
    add_card(slide, "引入三维几何构象", "融合原子坐标、键角、二面角与 SE(3) 等变网络，增强对空间结构特征的刻画。", 0.95, 1.65, 3.65, 1.50, TEAL, fill=WHITE, number="01")
    add_card(slide, "扩展多模态预训练规模", "构建更大规模、更高质量的分子图—序列—文本配对数据，提升跨任务泛化能力。", 4.85, 1.65, 3.65, 1.50, AMBER, fill=WHITE, number="02")
    add_card(slide, "干湿实验闭环验证", "将模型预测与真实实验反馈联动，形成主动学习式候选分子优化流程。", 8.75, 1.65, 3.65, 1.50, GREEN, fill=WHITE, number="03")
    add_image_contain(slide, "image/第四章/KGAMA及其变体散点结果.png", 1.05, 3.85, 4.10, 2.35,
                      caption="未来可继续提升回归任务的稳定性与校准能力", pad=0.12)
    add_text(slide, "谢谢各位老师\n请批评指正", 6.18, 4.20, 5.65, 0.95, size=28, color=WHITE,
             bold=True, align="center", margin=0, line_spacing=1.05)
    add_text(slide, "范帅尧 · 中国地质大学（北京）", 7.30, 5.35, 3.4, 0.26, size=12,
             color=MID, align="center", margin=0)
    add_footer(slide, 16, dark=True)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_deck()
