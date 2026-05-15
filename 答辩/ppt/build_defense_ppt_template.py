# -*- coding: utf-8 -*-
from pathlib import Path
import sys

sys.path.insert(0, str(Path(".codex_deps").resolve()))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE = next(
    (p for p in ROOT.glob("*.pptx") if p.name.startswith("ppt") and "-" in p.name and not p.name.startswith("~$")),
    ROOT / "ppt模板.pptx",
)
OUT = ROOT / "基于多模态学习的药物分子性质预测研究_答辩PPT.pptx"
ASSET_DIR = ROOT / "template_copy_assets"

SLIDE_W = 13.333333
SLIDE_H = 7.5

BLUE = "004E98"
BLUE_2 = "347AA7"
BLUE_3 = "7CA3C4"
DARK = "16324F"
PALE = "EAF3FA"
PAPER = "FFFFFF"
GRAY = "5B6773"
LIGHT_GRAY = "EEF2F5"
WHITE = "FFFFFF"
AMBER = "F5B54A"
RED = "B83232"

TITLE_FONT = "方正正中黑简体"
BODY_FONT = "微软雅黑 Light"
FALLBACK_FONT = "Microsoft YaHei"


def rgb(hex_color):
    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def inch(v):
    return Inches(v)


def delete_all_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def extract_template_assets():
    ASSET_DIR.mkdir(exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    picks = {
        "cover_bg": (1, 1),
        "logo": (2, 1),
    }
    outputs = {}
    for name, (slide_idx, pic_idx) in picks.items():
        slide = prs.slides[slide_idx - 1]
        pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
        image = pictures[pic_idx - 1].image
        path = ASSET_DIR / f"{name}.{image.ext}"
        path.write_bytes(image.blob)
        outputs[name] = path
    outputs["section_bg"] = outputs["cover_bg"]
    return outputs


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, radius=False):
    shp = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x, y, w, h=0, color=BLUE, width=1.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(max(h, 0.01)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    return shp


def add_text(slide, text, x, y, w, h, size=18, color=DARK, bold=False,
             align="left", valign="top", font=BODY_FONT, margin=0.04, line_spacing=1.05):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(margin)
    tf.margin_bottom = inch(margin)
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)
    for idx, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_picture_contain(slide, path, x, y, w, h, caption=None, border=True, pad=0.08, bg=WHITE):
    path = Path(path)
    if not path.is_absolute():
        path = ROOT / path
    if bg:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, bg, BLUE_3 if border else None)
    img = Image.open(path)
    iw, ih = img.size
    cap_h = 0.26 if caption else 0
    max_w = w - 2 * pad
    max_h = h - 2 * pad - cap_h
    ratio = min(max_w / iw, max_h / ih)
    dw, dh = iw * ratio, ih * ratio
    px = x + (w - dw) / 2
    py = y + pad + (max_h - dh) / 2
    slide.shapes.add_picture(str(path), inch(px), inch(py), width=inch(dw), height=inch(dh))
    if caption:
        add_text(slide, caption, x + 0.12, y + h - 0.27, w - 0.24, 0.18, size=8.4,
                 color=GRAY, align="center", margin=0)


def add_picture_cover(slide, path, x, y, w, h):
    path = Path(path)
    img = Image.open(path)
    iw, ih = img.size
    ratio = max(w / iw, h / ih)
    dw, dh = iw * ratio, ih * ratio
    px = x + (w - dw) / 2
    py = y + (h - dh) / 2
    slide.shapes.add_picture(str(path), inch(px), inch(py), width=inch(dw), height=inch(dh))


def add_header(slide, title, logo, part=None):
    add_line(slide, 0.95, 0.88, 10.95, 0.02, BLUE, 1)
    add_line(slide, 0.55, 0.53, 0.75, 0.035, BLUE, 1)
    add_line(slide, 0.77, 0.43, 0.55, 0.018, BLUE_2, 1)
    add_text(slide, title, 0.94, 0.32, 6.2, 0.45, size=20, color=BLUE, bold=True,
             font=TITLE_FONT, margin=0)
    if part:
        add_text(slide, part, 7.85, 0.37, 3.15, 0.30, size=9.5, color=BLUE_2,
                 align="right", margin=0)
    slide.shapes.add_picture(str(logo), inch(12.36), inch(0.34), width=inch(0.44), height=inch(0.50))


def add_footer(slide, page):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 7.18, SLIDE_W, 0.32, BLUE, None)
    add_text(slide, "艰苦朴素  求真务实", 0.0, 7.18, 2.42, 0.28, size=10, color=WHITE,
             align="center", valign="middle", margin=0)
    add_text(slide, f"{page:02d}", 12.45, 7.18, 0.42, 0.28, size=9, color=WHITE,
             align="right", valign="middle", margin=0)


def add_content_frame(slide, title, logo, page, part=None):
    set_bg(slide, WHITE)
    add_header(slide, title, logo, part)
    add_footer(slide, page)


def add_card(slide, title, body, x, y, w, h, accent=BLUE, fill=WHITE, num=None):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, BLUE_3)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.10, h, accent, None)
    if num:
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.22, y + 0.22, 0.46, 0.46, accent, None)
        add_text(slide, num, x + 0.22, y + 0.30, 0.46, 0.18, size=11, color=WHITE,
                 bold=True, align="center", margin=0)
        tx = x + 0.84
        tw = w - 1.10
    else:
        tx = x + 0.32
        tw = w - 0.58
    add_text(slide, title, tx, y + 0.22, tw, 0.30, size=14.5, color=DARK, bold=True,
             font=FALLBACK_FONT, margin=0)
    add_text(slide, body, x + 0.34, y + 0.68, w - 0.60, h - 0.78, size=10.5, color=GRAY,
             margin=0, line_spacing=1.08)


def add_metric(slide, value, label, x, y, w, h, color=BLUE):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, PALE, BLUE_3)
    add_text(slide, value, x + 0.05, y + 0.10, w - 0.10, 0.42, size=23, color=color,
             bold=True, align="center", margin=0, font=FALLBACK_FONT)
    add_text(slide, label, x + 0.10, y + 0.62, w - 0.20, h - 0.68, size=8.8, color=GRAY,
             align="center", margin=0)


def add_section_slide(prs, blank, bg, logo, number, title, subtitle):
    slide = prs.slides.add_slide(blank)
    add_picture_cover(slide, bg, 0, 0, SLIDE_W, SLIDE_H)
    slide.shapes.add_picture(str(logo), inch(12.36), inch(0.34), width=inch(0.44), height=inch(0.50))
    add_text(slide, f"PART {number}", 0.98, 1.05, 2.8, 0.38, size=17, color=WHITE,
             bold=True, font=FALLBACK_FONT, margin=0)
    add_text(slide, title, 0.98, 1.62, 7.2, 0.72, size=34, color=WHITE, bold=True,
             font=TITLE_FONT, margin=0)
    add_line(slide, 0.99, 2.58, 3.6, 0.035, WHITE)
    add_text(slide, subtitle, 1.00, 2.86, 7.8, 0.42, size=14, color="E7F0F8",
             font=BODY_FONT, margin=0)
    return slide


def add_flow(slide, items, x, y, w, h):
    gap = 0.18
    item_w = (w - gap * (len(items) - 1)) / len(items)
    colors = [BLUE, BLUE_2, BLUE_3, DARK]
    for i, (title, body) in enumerate(items):
        sx = x + i * (item_w + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, sx, y, item_w, h, colors[i % len(colors)], None)
        add_text(slide, f"{i + 1}", sx + 0.12, y + 0.12, 0.32, 0.20, size=11, color=WHITE,
                 bold=True, align="center", margin=0)
        add_text(slide, title, sx + 0.18, y + 0.48, item_w - 0.36, 0.30, size=12.2,
                 color=WHITE, bold=True, align="center", margin=0)
        add_text(slide, body, sx + 0.18, y + 0.86, item_w - 0.36, h - 0.94, size=8.8,
                 color="EAF3FA", align="center", margin=0, line_spacing=1.05)


def build():
    assets = extract_template_assets()
    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)
    blank = prs.slide_layouts[6]
    logo = assets["logo"]
    cover_bg = assets["cover_bg"]
    section_bg = assets["section_bg"]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_picture_cover(slide, cover_bg, 0, 0, SLIDE_W, SLIDE_H)
    add_text(slide, "基于多模态学习的药物分子性质预测研究", 0.20, 2.66, 12.90, 0.78,
             size=31, color=WHITE, bold=True, font=TITLE_FONT, align="center", margin=0)
    add_text(slide, "姓   名：范帅尧\n导   师：牛云云教授\n时   间：2026.05.13",
             7.86, 4.34, 4.20, 0.98, size=14, color=WHITE, font=BODY_FONT, margin=0, line_spacing=1.18)
    slide.shapes.add_picture(str(logo), inch(12.36), inch(0.34), width=inch(0.44), height=inch(0.50))

    # 2. Contents
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, -0.02, 0, 3.7, 7.5, BLUE, None)
    add_shape(slide, MSO_SHAPE.RIGHT_TRIANGLE, 2.80, 0, 1.45, 7.5, BLUE_2, None)
    add_text(slide, "目", 0.48, 1.10, 0.9, 0.90, size=42, color=WHITE, bold=True,
             font=TITLE_FONT, align="center", margin=0)
    add_text(slide, "录", 1.36, 1.98, 0.9, 0.90, size=42, color=WHITE, bold=True,
             font=TITLE_FONT, align="center", margin=0)
    add_text(slide, "CONTENTS", 0.40, 3.55, 2.35, 0.34, size=17, color="D9EAF7",
             font="Aptos", margin=0)
    slide.shapes.add_picture(str(logo), inch(12.36), inch(0.34), width=inch(0.44), height=inch(0.50))
    contents = [
        ("01", "研究背景与意义", "药物发现早筛需求与多模态分子性质预测价值"),
        ("02", "研究现状", "单模态、多模态方法进展及其主要不足"),
        ("03", "研究内容", "CoT-CMP 与 KGAMA 两项核心方法及实验验证"),
        ("04", "结论与展望", "研究贡献、局限性与未来工作方向"),
    ]
    for i, (num, title, desc) in enumerate(contents):
        y = 1.18 + i * 1.28
        add_text(slide, f"{num} /", 5.05, y, 0.78, 0.26, size=15, color=BLUE,
                 bold=True, font=FALLBACK_FONT, margin=0)
        add_text(slide, title, 5.88, y - 0.02, 2.8, 0.30, size=18, color=DARK,
                 bold=True, font=FALLBACK_FONT, margin=0)
        add_text(slide, desc, 5.88, y + 0.44, 4.7, 0.24, size=10.4, color=GRAY, margin=0)
        add_line(slide, 5.88, y + 0.34, 4.55, 0.012, BLUE_3)

    # 3. Section 1
    add_section_slide(prs, blank, section_bg, logo, "01", "研究背景与意义", "面向药物早期筛选的分子性质预测需求")

    # 4. Background
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "研究背景：药物早筛对性质预测提出高要求", logo, 4, "第一部分 研究背景与意义")
    add_text(slide, "分子性质预测旨在建立分子结构与理化属性、生化性质之间的映射关系，是现代药物发现早期筛选中的重要环节。",
             0.95, 1.30, 5.95, 0.62, size=16, color=DARK, bold=True, font=FALLBACK_FONT, margin=0)
    add_card(slide, "研发流程复杂", "创新药研发周期长、成本高，后期失败风险会显著放大前期筛选误差。", 0.95, 2.25, 2.75, 1.04, BLUE, num="1")
    add_card(slide, "实验验证昂贵", "候选分子数量巨大，湿实验难以覆盖全部化学空间。", 3.98, 2.25, 2.75, 1.04, BLUE_2, num="2")
    add_card(slide, "AI 可前置风险", "通过计算模型预测 ADMET、毒性和理化性质，辅助候选分子优先级排序。", 7.00, 2.25, 2.75, 1.04, BLUE_3, num="3")
    add_picture_contain(slide, "image/第二章相关理论基础/图2-1 阿司匹林的多模态化学表征.png",
                        0.95, 4.02, 5.75, 2.10, "同一分子可从图结构、SMILES、指纹和文本语义等角度表征", pad=0.10)
    add_picture_contain(slide, "image/第二章相关理论基础/图2-2 消息传递机制示意图.png",
                        7.02, 4.02, 5.35, 2.10, "图神经网络为分子结构建模提供基础", pad=0.10)

    # 5. Meaning
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "研究意义：提升准确性、可解释性与泛化能力", logo, 5, "第一部分 研究背景与意义")
    add_flow(slide, [
        ("候选分子", "海量化合物库"),
        ("多模态表征", "结构、序列、指纹、文本"),
        ("性质预测", "分类与回归任务"),
        ("早期决策", "风险筛除与优先级排序"),
    ], 0.95, 1.28, 11.35, 1.32)
    add_card(slide, "准确性", "充分利用分子多视角信息，弥补单一模态对复杂结构与药理语义表达不足的问题。", 0.95, 3.28, 3.48, 1.22, BLUE, num="01")
    add_card(slide, "可解释性", "将大语言模型的推理能力转化为显式逻辑文本，使预测过程更容易被理解与审计。", 4.90, 3.28, 3.48, 1.22, BLUE_2, num="02")
    add_card(slide, "泛化能力", "通过知识增强与图中心对齐降低文本噪声和模态冲突，提升跨数据集稳定表现。", 8.85, 3.28, 3.48, 1.22, BLUE_3, num="03")
    add_text(slide, "本文围绕多模态分子性质预测，提出两种方法：CoT-CMP 侧重可解释图文融合，KGAMA 侧重知识增强与稳健对齐。",
             1.15, 5.52, 10.85, 0.50, size=16.5, color=BLUE, bold=True, font=FALLBACK_FONT, align="center", margin=0)

    # 6. Section 2
    add_section_slide(prs, blank, section_bg, logo, "02", "研究现状", "分子性质预测方法进展与当前不足")

    # 7. Research status
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "研究现状：从单模态建模到多模态融合", logo, 7, "第二部分 研究现状")
    add_picture_contain(slide, "image/第二章相关理论基础/图2-3 图卷积神经网络结构示意图.png",
                        0.95, 1.22, 3.55, 2.12, "图神经网络：结构拓扑建模", pad=0.10)
    add_picture_contain(slide, "image/第二章相关理论基础/图2-4 Transformer架构.png",
                        4.88, 1.22, 3.55, 2.12, "Transformer：序列与文本建模", pad=0.10)
    add_picture_contain(slide, "image/第二章相关理论基础/LoRA技术原理示意图.png",
                        8.82, 1.22, 3.55, 2.12, "LoRA：大模型高效适配", pad=0.10)
    add_card(slide, "单模态方法", "基于分子图、SMILES、指纹或理化描述分别建模，结构清晰但信息来源有限。", 0.95, 3.85, 3.50, 1.28, BLUE)
    add_card(slide, "多模态方法", "将图、序列、文本等信息联合学习，能补充分子结构与药理语义之间的信息缺口。", 4.92, 3.85, 3.50, 1.28, BLUE_2)
    add_card(slide, "大模型方法", "LLM 具备领域知识和文本生成能力，但直接用于性质预测时仍需解决幻觉与解释约束问题。", 8.88, 3.85, 3.50, 1.28, BLUE_3)
    add_text(slide, "总体趋势：从“只看结构”逐步走向“结构 + 语义 + 先验知识”的综合表征学习。",
             1.0, 5.70, 11.3, 0.38, size=15, color=BLUE, bold=True, align="center", margin=0)

    # 8. Gaps
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "当前研究不足", logo, 8, "第二部分 研究现状")
    add_card(slide, "拓扑感知受限与信息回流", "传统图神经网络多依赖局部节点聚合，化学键动态交互建模不足；深层传播易出现信息冗余。", 0.95, 1.32, 3.55, 1.62, BLUE, num="01")
    add_card(slide, "LLM 直接预测难以解释", "面对复杂 ADMET 任务时，端到端文本映射缺少中间药理逻辑约束，容易产生幻觉。", 4.88, 1.32, 3.55, 1.62, BLUE_2, num="02")
    add_card(slide, "文本噪声与模态对齐困难", "公共医学文本存在冗余和质量差异，粗粒度对比学习可能引入语义冲突和负迁移。", 8.82, 1.32, 3.55, 1.62, BLUE_3, num="03")
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.12, 4.12, 5.20, 1.12, PALE, BLUE_3)
    add_text(slide, "CoT-CMP", 1.45, 4.34, 1.45, 0.28, size=17, color=BLUE, bold=True, font=FALLBACK_FONT, margin=0)
    add_text(slide, "通过思维链增强文本可解释性，通过 CMPNN 强化图结构交互建模。", 1.45, 4.72, 4.35, 0.26, size=11.2, color=GRAY, margin=0)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.02, 4.12, 5.20, 1.12, PALE, BLUE_3)
    add_text(slide, "KGAMA", 7.35, 4.34, 1.45, 0.28, size=17, color=BLUE, bold=True, font=FALLBACK_FONT, margin=0)
    add_text(slide, "通过知识重写提高文本质量，通过图中心对齐与不确定性加权稳定融合。", 7.35, 4.72, 4.35, 0.26, size=11.2, color=GRAY, margin=0)

    # 9. Section 3
    add_section_slide(prs, blank, section_bg, logo, "03", "研究内容", "两项多模态分子性质预测方法设计与实验验证")

    # 10. Research content overview
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "研究内容整体框架", logo, 10, "第三部分 研究内容")
    add_card(slide, "研究内容一：CoT-CMP", "面向 TDC ADMET 基准组，构建思维链增强与通信式消息传递结合的图文双塔预测框架。", 0.95, 1.35, 5.45, 1.38, BLUE, num="01")
    add_card(slide, "研究内容二：KGAMA", "面向 MoleculeNet 基准组，构建知识增强、图中心对齐和四模态融合的分子性质预测框架。", 6.92, 1.35, 5.45, 1.38, BLUE_2, num="02")
    add_picture_contain(slide, "image/第三章/图3-1CoTCMP架构图.png",
                        0.95, 3.12, 5.45, 2.25, "CoT-CMP：文本推理 + 图结构编码 + 交叉注意力", pad=0.10)
    add_picture_contain(slide, "image/第四章/KGAMA模型架构.png",
                        6.92, 3.12, 5.45, 2.25, "KGAMA：图、序列、指纹、文本四模态对齐", pad=0.10)
    add_text(slide, "核心思想：不是简单堆叠模态，而是让模态具备知识约束、可靠对齐和任务自适应贡献。",
             1.10, 5.95, 11.0, 0.34, size=15, color=BLUE, bold=True, align="center", margin=0)

    # 11. CoT-CMP method
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "方法一：CoT-CMP 双塔预测框架", logo, 11, "第三部分 研究内容")
    add_picture_contain(slide, "image/第三章/图3-1CoTCMP架构图.png",
                        0.95, 1.18, 7.35, 3.10, "图 3-1 CoT-CMP 架构图", pad=0.10)
    add_card(slide, "文本编码端", "LLaMA-3.1-8B + LoRA + CoT 提示，生成结构识别、理化分析、性质推断的显式推理文本。", 8.65, 1.24, 3.72, 1.22, BLUE)
    add_card(slide, "图结构编码端", "采用 CMPNN 建模原子与化学键之间的通信关系，增强局部官能团与长程依赖表征。", 8.65, 2.77, 3.72, 1.22, BLUE_2)
    add_card(slide, "融合预测端", "通过交叉注意力实现图特征与文本特征的自适应融合，服务分类与回归任务。", 8.65, 4.30, 3.72, 1.05, BLUE_3)
    add_picture_contain(slide, "table/第三章/提示词模板.png",
                        0.95, 4.62, 3.55, 1.62, "CoT 提示词模板", pad=0.08)
    add_picture_contain(slide, "image/第二章相关理论基础/图2-6 输入-输出提示与思维链提示对比示意图.png",
                        4.78, 4.62, 3.52, 1.62, "思维链提示对比", pad=0.08)

    # 12. CoT-CMP results
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "CoT-CMP 实验结果", logo, 12, "第三部分 研究内容")
    add_picture_contain(slide, "table/第三章/CoT-CMP在ADMET基准组上的对比实验结果.png",
                        0.95, 1.18, 6.28, 4.92, "表 3-3 CoT-CMP 在 ADMET 基准组上的对比实验结果", pad=0.08)
    add_metric(slide, "17/22", "ADMET 数据集取得最优", 7.68, 1.30, 1.75, 0.88, BLUE)
    add_metric(slide, "5/22", "取得次优结果", 9.58, 1.30, 1.55, 0.88, BLUE_2)
    add_metric(slide, "4/4", "毒性任务全面领先", 11.25, 1.30, 1.18, 0.88, RED)
    add_card(slide, "对比实验结论", "CoT-CMP 在 22 个 ADMET 评测集上获得 17 个最优、5 个次优；在 LD50、hERG、Ames、DILI 四项毒性预测任务上全面领先。", 7.68, 2.58, 4.75, 1.22, BLUE)
    add_card(slide, "性能提升原因", "CoT 增强文本描述与交叉注意力融合能够更有效聚合长程依赖和局部官能团特征，形成更鲁棒的分子表示。", 7.68, 4.18, 4.75, 1.22, BLUE_2)

    # 13. CoT-CMP ablation
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "CoT-CMP 消融与可解释性分析", logo, 13, "第三部分 研究内容")
    add_picture_contain(slide, "table/第三章/CoT-CMP在ADMET基准组上部分分类任务消融实验结果.png",
                        0.95, 1.20, 5.42, 1.58, "分类任务消融结果", pad=0.08)
    add_picture_contain(slide, "image/第三章/氯磺丙脲在DILI任务中的双模态可视化分析结果.png",
                        0.95, 3.12, 5.42, 2.92, "DILI 任务中的双模态可视化案例", pad=0.08)
    add_picture_contain(slide, "image/第三章/CMPNN模型层数对分类任务性能的影响.png",
                        6.78, 1.20, 5.55, 2.36, "CMPNN 层数敏感性实验", pad=0.08)
    add_card(slide, "消融结论", "移除 LLM、CoT、LoRA 或交叉注意力均会造成不同程度下降，说明文本增强、任务适配和融合模块存在协同增益。", 6.78, 3.95, 5.55, 0.98, BLUE)
    add_card(slide, "解释结论", "模型关注对氯苯基、磺酰脲主骨架以及 DILI、hepatotoxicity 等关键词，与毒理学判断直觉一致。", 6.78, 5.18, 5.55, 0.86, BLUE_2)

    # 14. KGAMA method
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "方法二：KGAMA 知识增强与图对齐框架", logo, 14, "第三部分 研究内容")
    add_picture_contain(slide, "image/第四章/KGAMA模型架构.png",
                        0.95, 1.15, 7.10, 4.70, "图 4-1 KGAMA 模型架构", pad=0.08)
    add_card(slide, "四模态表征", "联合分子图、SMILES 序列、分子指纹和自然语言文本，构建更完整的分子表示体系。", 8.42, 1.18, 3.90, 1.05, BLUE)
    add_card(slide, "知识重写", "利用 PubChem 原始文本和 RDKit 计算的 11 项理化指标，经 LLM 去噪重写后输入 SciBERT。", 8.42, 2.54, 3.90, 1.05, BLUE_2)
    add_card(slide, "图中心对齐", "以图表征作为稳定语义锚点，并用同方差不确定性权重降低噪声模态干扰。", 8.42, 3.90, 3.90, 1.05, BLUE_3)
    add_card(slide, "自适应融合", "通过交叉注意力根据具体任务动态调配不同模态贡献，而非静态拼接。", 8.42, 5.26, 3.90, 0.78, DARK)

    # 15. KGAMA knowledge rewrite
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "KGAMA 知识重写与模态构建", logo, 15, "第三部分 研究内容")
    add_flow(slide, [
        ("PubChem", "原始分子描述"),
        ("RDKit", "理化指标注入"),
        ("LLM 重写", "去噪与知识融合"),
        ("SciBERT", "高密度语义编码"),
    ], 0.95, 1.22, 6.15, 1.22)
    add_picture_contain(slide, "image/第四章/药物分子原始检索文本与重写文本的对比示例.png",
                        7.55, 1.08, 4.75, 5.20, "药物分子原始检索文本与重写文本对比示例", pad=0.08)
    add_picture_contain(slide, "table/第四章/三种典型分子指纹特征的比较分析.png",
                        0.95, 3.04, 2.90, 1.15, "三类分子指纹特征", pad=0.06)
    add_picture_contain(slide, "table/第四章/分子理化性质指标及其含义.png",
                        4.18, 3.04, 2.92, 1.15, "理化性质指标", pad=0.06)
    add_card(slide, "知识增强价值", "将冗余公共医学文本转换为结构化、化学约束更强的专家语义，提升文本模态可用性。", 0.95, 4.62, 2.90, 1.15, BLUE)
    add_card(slide, "指纹补充价值", "AtomPairs、MACCS、Morgan 从不同角度捕获子结构统计信息，弥补图编码器可能遗漏的全局模式。", 4.18, 4.62, 2.92, 1.15, BLUE_2)

    # 16. KGAMA results
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "KGAMA 实验结果与可解释性", logo, 16, "第三部分 研究内容")
    add_picture_contain(slide, "table/第四章/KGAMA在MoleculeNet分类任务上对比实验结果.png",
                        0.95, 1.12, 5.55, 3.10, "KGAMA 在 MoleculeNet 分类任务上的对比实验结果", pad=0.08)
    add_picture_contain(slide, "image/第四章/不同分子模态在各下游任务中的注意力权重分布.png",
                        6.90, 1.12, 5.42, 3.10, "不同分子模态在各下游任务中的注意力权重分布", pad=0.08)
    add_metric(slide, "0.939", "ClinTox AUROC", 0.95, 4.60, 1.55, 0.76, RED)
    add_metric(slide, "0.930", "BBBP AUROC", 2.72, 4.60, 1.55, 0.76, BLUE)
    add_metric(slide, "0.831", "ESOL RMSE", 4.50, 4.60, 1.55, 0.76, BLUE_2)
    add_card(slide, "实验结论", "KGAMA 在 MoleculeNet 分类与回归任务上均取得较好综合表现，验证知识增强和图中心对齐策略有效。", 6.90, 4.50, 5.42, 0.82, BLUE)
    add_card(slide, "注意力解释", "物理性质任务更依赖 Graph / Fingerprint；毒性与副作用任务中文本模态权重显著提升。", 6.90, 5.55, 5.42, 0.70, BLUE_2)

    # 17. Section 4
    add_section_slide(prs, blank, section_bg, logo, "04", "结论与展望", "研究贡献总结与未来工作方向")

    # 18. Conclusion and outlook
    slide = prs.slides.add_slide(blank)
    add_content_frame(slide, "结论与展望", logo, 18, "第四部分 结论与展望")
    add_card(slide, "结论一：CoT-CMP", "提出思维链增强与通信式图交互的图文双塔框架，在 ADMET 基准组上取得较优预测结果，并提供可解释线索。", 0.95, 1.25, 5.45, 1.36, BLUE, num="01")
    add_card(slide, "结论二：KGAMA", "提出知识增强与图中心对齐的四模态框架，在 MoleculeNet 基准组上验证了知识重写和自适应融合的有效性。", 6.92, 1.25, 5.45, 1.36, BLUE_2, num="02")
    add_card(slide, "展望一：引入三维几何构象", "进一步融合原子坐标、键角、二面角及 SE(3) 等变模型，提升空间结构建模能力。", 0.95, 3.18, 3.50, 1.22, BLUE)
    add_card(slide, "展望二：扩展预训练数据规模", "构建更大规模、高质量的图、序列、文本配对数据，增强泛化能力。", 4.92, 3.18, 3.50, 1.22, BLUE_2)
    add_card(slide, "展望三：干湿实验闭环验证", "将模型预测与真实实验反馈联动，服务候选分子的持续优化。", 8.88, 3.18, 3.50, 1.22, BLUE_3)
    add_text(slide, "谢谢各位老师，请批评指正", 2.35, 5.62, 8.65, 0.54, size=28,
             color=BLUE, bold=True, font=TITLE_FONT, align="center", margin=0)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
